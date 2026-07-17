import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Theme Colors
MIDNIGHT_GREEN = RGBColor(15, 42, 29)     # #0F2A1D
REFINED_GOLD = RGBColor(212, 162, 76)     # #D4A24C
PEACH_GOLD = RGBColor(243, 192, 107)      # #F3C06B
OBSIDIAN_BLACK = RGBColor(18, 20, 18)     # #121412
CARD_SURFACE = RGBColor(20, 42, 34)       # #142A22
WHITE = RGBColor(255, 255, 255)
WHITE_DIM = RGBColor(200, 200, 200)

def apply_slide_background(slide, color):
    """Sets the background color of a slide using a solid fill rectangle."""
    left = top = Inches(0)
    width = Inches(13.33)
    height = Inches(7.5)
    rect = slide.shapes.add_shape(
        1, left, top, width, height # 1 is shape type for Rectangle
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background() # No border
    # Send to back
    slide.shapes._spTree.remove(rect._element)
    slide.shapes._spTree.insert(2, rect._element)

def add_header(slide, title_text, category_text="AL SAFA AR DINING"):
    """Adds a standard premium header to the slide."""
    # Category / Brand small text
    cat_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(10), Inches(0.3))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.size = Pt(11)
    p_cat.font.bold = True
    p_cat.font.color.rgb = REFINED_GOLD
    p_cat.font.name = 'Outfit'
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.65), Inches(10), Inches(0.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.size = Pt(32)
    p_title.font.bold = True
    p_title.font.color.rgb = WHITE
    p_title.font.name = 'Outfit'

def add_bullet_list(slide, items, left, top, width, height, font_size=16):
    """Adds a bulleted list to a slide inside a textbox."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = WHITE_DIM
        p.font.name = 'Inter'
        p.space_after = Pt(8)
        p.level = 0

def add_card(slide, title, bullets, left, top, width, height):
    """Adds an elegant container card with text content."""
    # Background card shape
    rect = slide.shapes.add_shape(1, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = CARD_SURFACE
    rect.line.color.rgb = REFINED_GOLD
    rect.line.width = Pt(1)
    
    # Card Title
    title_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.1), width - Inches(0.4), Inches(0.5))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title.upper()
    p_title.font.size = Pt(14)
    p_title.font.bold = True
    p_title.font.color.rgb = PEACH_GOLD
    p_title.font.name = 'Outfit'
    
    # Card bullets
    add_bullet_list(slide, bullets, left + Inches(0.2), top + Inches(0.5), width - Inches(0.4), height - Inches(0.6), font_size=13)

def generate_proposal():
    print("Generating Task 1: Proposal Slide Deck...")
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    # ------------------ SLIDE 1: Title Slide ------------------
    slide = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide, MIDNIGHT_GREEN)
    
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.2), Inches(11.8), Inches(2.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "AL SAFA AR DINING"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = REFINED_GOLD
    p.font.name = 'Outfit'
    
    p2 = tf.add_paragraph()
    p2.text = "Task 1: Interactive AR Digital Menu Proposal (80-Item Database)"
    p2.font.size = Pt(24)
    p2.font.color.rgb = WHITE
    p2.font.name = 'Outfit'
    p2.space_before = Pt(12)
    
    sub_box = slide.shapes.add_textbox(Inches(0.75), Inches(4.8), Inches(10), Inches(1.5))
    tf_sub = sub_box.text_frame
    p_sub1 = tf_sub.paragraphs[0]
    p_sub1.text = "Course: Introduction to Multimedia Technology (BMD12083)\nPrepared for: Raffles University Johor Bahru\nSubmitted by: Shanjaay (Student ID: 202420058)"
    p_sub1.font.size = Pt(14)
    p_sub1.font.color.rgb = WHITE_DIM
    p_sub1.font.name = 'Inter'

    img_path = 'assets/restaurant_bg.png'
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(8.5), Inches(2.2), width=Inches(4.0))

    # ------------------ SLIDE 2: Project Introduction ------------------
    slide = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide, OBSIDIAN_BLACK)
    add_header(slide, "Project Introduction")
    
    bullets_intro = [
        "Mamak Food Culture: Traditional menus lack interactive visual contexts and customization transparency, leading to orders that fail to match user customized tastes.",
        "The Solution: Al Safa AR Dining is a responsive web application offering a massive digital menu database containing all 80 items from Mamak_Menu_Malaysia.docx.",
        "Real-Time Overlays: Allows customers to adjust specifications (egg yolk paste, cheese slices, chili rings, flooded curry) and watch changes update instantly.",
        "Quality Aesthetics: A luxury obsidian black, midnight green, and gold accents design system optimized for modern ambient dining layouts."
    ]
    add_bullet_list(slide, bullets_intro, Inches(0.75), Inches(1.8), Inches(7.5), Inches(4.8))
    
    img_path = 'assets/nasi_lemak.png'
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(8.8), Inches(1.8), width=Inches(3.8))

    # ------------------ SLIDE 3: 80-Item Menu Database ------------------
    slide = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide, OBSIDIAN_BLACK)
    add_header(slide, "Comprehensive 80-Item Menu Database")
    
    bullets_mains = [
        "Full Menu Integration: Loaded all 80 menu items from the official Mamak Menu database, ranging from RM1.95 (Teh O) up to RM16.90 (Nasi Goreng Special).",
        "Categorized Navigation: Features 8 specific category filters: Roti/Flatbreads (15), Nasi Goreng (16), Noodles (9), Nasi Dishes (3), Sides/Lauk (5), Snacks/Desserts (6), Hot Drinks (10), Cold Drinks (11), and Specialty Drinks (5).",
        "Visual Assets Mapping: High-fidelity renders act as structural templates to dynamically represent all 80 distinct items."
    ]
    add_bullet_list(slide, bullets_mains, Inches(0.75), Inches(1.8), Inches(7.5), Inches(4.8))
    
    img_path = 'assets/nasi_kandar.png'
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(8.8), Inches(1.8), width=Inches(3.8))

    # ------------------ SLIDE 4: Category-Specific Visual Customizers ------------------
    slide = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide, OBSIDIAN_BLACK)
    add_header(slide, "Real-Time Visual Customizers (AR Canvas)")
    
    # 2 Cards for Customizers
    bullets_drinks = [
        "Roti Flatbreads: Programmatically overlays melted cheese grids, sunny egg yolks, or white condensed milk spirals.",
        "Fried Rice & Noodles: Toggle spiciness (Mild, Medium, Extra Hot) to scatter red chili rings and apply red glow filters.",
        "Curry Flood (Banjir): Adjust curry gravy height on the food plate visual using custom sliders."
    ]
    add_card(slide, "Food Customizers (AR Overlays)", bullets_drinks, Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.5))
    
    bullets_3d = [
        "Drinks Customizer: Ice level adjustments spawn transparent ice cubes; milk/rose cordial levels adjust opacity and tea/pink milk hues.",
        "Milo Dinosaur: Adjust Milo powder slider to dynamically expand or scale the brown cocoa heap representation.",
        "Zoom & Rotate: Perform 360-degree rotations and up to 1.8x scale transforms on 3D plate models."
    ]
    add_card(slide, "Drink Sliders & Zoom Dials", bullets_3d, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.5))

    # ------------------ SLIDE 5: Technology & Software ------------------
    slide = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide, OBSIDIAN_BLACK)
    add_header(slide, "Technology & Software Stack")
    
    bullets_tech_stack = [
        "Application Framework: Flutter (Dart) for high-performance cross-platform web and mobile compilations.",
        "Design & Modeling: Figma for high-fidelity UI layout vectors, and Blender for 3D modeling of foods.",
        "Asset Compilation: Adobe Photoshop for texture editing; image generation and processing.",
        "Interactive Engine: Programmatic transformation matrices and custom overlay painters simulating spatial camera views."
    ]
    add_bullet_list(slide, bullets_tech_stack, Inches(0.75), Inches(1.8), Inches(7.5), Inches(4.8))
    
    img_path = 'assets/milo_dinosaur.png'
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(8.8), Inches(1.8), width=Inches(3.8))

    # ------------------ SLIDE 6: Cost & Budgeting ------------------
    slide = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide, OBSIDIAN_BLACK)
    add_header(slide, "Project Cost & Budget Estimation")
    
    bullets_costs = [
        "Software Licensing: Adobe Creative Cloud & Blender (RM 250 / Month - Student License)",
        "Hardware Resources: High-Performance GPU rendering workstation (RM 4,500 - One-time)",
        "Assets & Hosting: Web deployment server and cloud assets storage (RM 120 / Month)",
        "Printed Marketing: NFC-integrated QR wooden table stands (RM 400 - One-time)",
        "Estimated Total Project Budget: RM 5,770 (Initial setup and first 3 months operating costs)"
    ]
    add_bullet_list(slide, bullets_costs, Inches(0.75), Inches(1.8), Inches(11.8), Inches(4.5))

    # ------------------ SLIDE 7: Project Timeline ------------------
    slide = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide, OBSIDIAN_BLACK)
    add_header(slide, "Project Development Timeline")
    
    # Adding a table for timeline
    rows, cols = 5, 3
    left, top, width, height = Inches(0.75), Inches(1.8), Inches(11.8), Inches(4.5)
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    # Header cells
    table.cell(0, 0).text = "PHASE / TASK"
    table.cell(0, 1).text = "DURATION"
    table.cell(0, 2).text = "DELIVERABLES"
    
    timeline_data = [
        ("Phase 1: Planning & Research", "Weeks 1 - 3", "Project proposal, storyboard drafts, user journey maps."),
        ("Phase 2: Asset Production", "Weeks 4 - 7", "3D food models, background plates, photography styling."),
        ("Phase 3: App Development", "Weeks 8 - 11", "Flutter structure, interactive AR simulator logic, UI cards."),
        ("Phase 4: Testing & Presentation", "Weeks 12 - 13", "User testing, design journal slides, final project reflection.")
    ]
    
    for row_idx, data in enumerate(timeline_data, start=1):
        for col_idx, text in enumerate(data):
            table.cell(row_idx, col_idx).text = text
            
    # Set text colors and styling for table cells
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.name = 'Inter'
                if r == 0:
                    p.font.bold = True
                    p.font.color.rgb = REFINED_GOLD
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = MIDNIGHT_GREEN
                else:
                    p.font.color.rgb = WHITE_DIM
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = CARD_SURFACE

    # ------------------ SLIDE 8: Sketches & Storyboard ------------------
    slide = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide, OBSIDIAN_BLACK)
    add_header(slide, "Sketches & User Storyboard")
    
    bullets_story = [
        "1. Scan QR Stand: Customer scans the wooden table QR stand, which launches the camera scanner and unlocks Table 12.",
        "2. Browse Menu: User is presented with a premium layout containing 80 categorized items.",
        "3. Detail Selection: Patron selects a dish (e.g. Roti Canai), chooses toppings (egg/cheese) or spiciness.",
        "4. Interactive Customizer: Taps 'VIEW IN AR' to adjust sliders (gravy banjir/sweetness) and verify visual representation changes in real-time.",
        "5. Checkout & Confirmed: Order is placed directly to Table 12, estimating 5-10 minutes preparation."
    ]
    add_bullet_list(slide, bullets_story, Inches(0.75), Inches(1.8), Inches(7.5), Inches(4.8), font_size=13)
    
    img_path = 'assets/qr_stand.png'
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(8.8), Inches(1.8), width=Inches(3.8))

    # Save presentation
    prs.save("Task_1_Proposal.pptx")
    print("Task 1: Proposal Slide Deck saved successfully!")


def generate_design_journal():
    print("Generating Task 2: Design Journal Slide Deck...")
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    # ------------------ SLIDE 1: Title Slide ------------------
    slide = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide, MIDNIGHT_GREEN)
    
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.2), Inches(11.8), Inches(2.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "AL SAFA AR DINING"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = REFINED_GOLD
    p.font.name = 'Outfit'
    
    p2 = tf.add_paragraph()
    p2.text = "Task 2: Design Journal & Production Process (80-Item Menu)"
    p2.font.size = Pt(24)
    p2.font.color.rgb = WHITE
    p2.font.name = 'Outfit'
    p2.space_before = Pt(12)
    
    sub_box = slide.shapes.add_textbox(Inches(0.75), Inches(4.8), Inches(10), Inches(1.5))
    tf_sub = sub_box.text_frame
    p_sub1 = tf_sub.paragraphs[0]
    p_sub1.text = "Prepared by: Shanjaay (Student ID: 202420058)\nAssignment 3 - Production & Prototype Showcase"
    p_sub1.font.size = Pt(14)
    p_sub1.font.color.rgb = WHITE_DIM
    p_sub1.font.name = 'Inter'
    
    img_path = 'assets/satay.png'
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(8.5), Inches(2.2), width=Inches(4.0))

    # ------------------ SLIDE 2: Brainstorming & Research ------------------
    slide = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide, OBSIDIAN_BLACK)
    add_header(slide, "Stage 1: Brainstorming & User Pain Points")
    
    bullets_brain = [
        "Mamak restaurants are often packed, causing order delays or misunderstandings regarding customizations.",
        "Portion sizes and ingredients (like standard vs sharing curry) can be ambiguous, leading to food waste or under-ordering.",
        "We brainstormed an AR tool allowing patrons to see exactly what they are ordering directly on their tables, modifying parameters beforehand.",
        "The project evolved from simple 2D icons to a high-fidelity, interactive, and spatial visual components database representing 80 menu items."
    ]
    add_bullet_list(slide, bullets_brain, Inches(0.75), Inches(1.8), Inches(11.8), Inches(4.5))

    # ------------------ SLIDE 3: Moodboard & Design Tokens ------------------
    slide = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide, OBSIDIAN_BLACK)
    add_header(slide, "Stage 2: Style Guide & Moodboard")
    
    bullets_mood = [
        "Primary Green (#0F2A1D): Midnight forest tone, representing Islamic-Mamak architectural heritage.",
        "Accent Gold (#D4A24C): Luxurious brass glow, highlighting premium prices, callout pins, and active buttons.",
        "Surface Obsidian (#121412): Modern, deep black background ensuring perfect visibility for translucent widgets.",
        "Glassmorphism: Opacity of 70% with 20px backdrop blurs. Inner 0.5px white borders catch the ambient light.",
        "Typography: Outfit for headings; Inter for structural readability in mobile AR canvases."
    ]
    add_bullet_list(slide, bullets_mood, Inches(0.75), Inches(1.8), Inches(7.5), Inches(4.8), font_size=14)
    
    img_path = 'assets/teh_tarik.png'
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(8.8), Inches(1.8), width=Inches(3.8))

    # ------------------ SLIDE 4: Behind the Scenes - Asset Creation ------------------
    slide = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide, OBSIDIAN_BLACK)
    add_header(slide, "Stage 3: Behind the Scenes (Assets Creation)")
    
    bullets_bts = [
        "1. 3D Food Modeling: Using Blender to sculpt Roti Canai folds, Roti Tisu cone cylinders, Mee Goreng noodles, and Teh Tarik froth layers.",
        "2. Material & Shading: Applying procedural shaders with condensation and light refraction settings to capture food steam and gloss.",
        "3. Multi-angle Rendering: Exporting model angles to support rotation inside the Flutter canvas.",
        "4. Image Optimization: Converting large renders into web-optimized WebP assets (~30% size reduction) to ensure low-latency loading."
    ]
    add_bullet_list(slide, bullets_bts, Inches(0.75), Inches(1.8), Inches(7.5), Inches(4.8), font_size=14)
    
    img_path = 'assets/mee_goreng.png'
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(8.8), Inches(1.8), width=Inches(3.8))

    # ------------------ SLIDE 5: Interactive AR Prototype ------------------
    slide = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide, OBSIDIAN_BLACK)
    add_header(slide, "Stage 4: Interactive AR Prototype")
    
    bullets_proto = [
        "QR Table Stand Scan: Opens a mock camera feed containing a golden boundary. Successfully matches the QR code to lock Table 12 context.",
        "Visual Customizer Overlays: Sliders and toggles programmatically draw melted cheese grids, egg yolk centers, white condensed milk swirls, and red chili rings.",
        "Curry Flood (Banjir): Slides adjust gravy layers on the food visual dynamically.",
        "Nasi Lemak Callouts: Callout nodes overlay the coconut rice, sambal, and fried chicken, providing descriptions of slow-cooked recipes on tap."
    ]
    add_bullet_list(slide, bullets_proto, Inches(0.75), Inches(1.8), Inches(7.5), Inches(4.8), font_size=14)
    
    img_path = 'assets/qr_stand.png'
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(8.8), Inches(1.8), width=Inches(3.8))

    # ------------------ SLIDE 6: Production & Code Architecture ------------------
    slide = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide, OBSIDIAN_BLACK)
    add_header(slide, "Stage 5: Production & Flutter Implementation")
    
    bullets_code = [
        "1. State Management: Provider package coordinates cart states (quantities, notes, portion modifiers) globally across pages.",
        "2. Custom Painters: The AR camera overlay and roti/milk toppings utilize CustomPainters to render vector outlines and grids at runtime.",
        "3. Performance Tree-Shaking: Optimized compile-time icon shaking, asset configurations, and font caching reduces web bundles.",
        "4. Glassmorphism Widgets: BackdropFilter and BoxDecoration combined to create high-end translucent glass cards."
    ]
    add_bullet_list(slide, bullets_code, Inches(0.75), Inches(1.8), Inches(11.8), Inches(4.5))

    # Save presentation
    prs.save("Task_2_Design_Journal.pptx")
    print("Task 2: Design Journal Slide Deck saved successfully!")

if __name__ == '__main__':
    generate_proposal()
    generate_design_journal()
    print("All slides generated successfully!")
